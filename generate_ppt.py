from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_pres():
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Colors
    c_deep_slate = RGBColor(26, 28, 31)
    c_white = RGBColor(255, 255, 255)
    c_gold = RGBColor(211, 175, 121)
    c_blue = RGBColor(40, 119, 188)
    c_purple = RGBColor(84, 84, 163)
    c_black = RGBColor(0, 0, 0)
    c_grey = RGBColor(100, 100, 100)

    # Margins (approximate conversions for 13.333 x 7.5, roughly 1 inch = 72 pt)
    # 80px = ~ 1.11 inches
    safe_margin = Inches(1.11)

    # 1. Cover
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_deep_slate
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "AI 기반의 프로덕트 설계: 앱 제작 프로세스의 재설계"
    p.font.size = Pt(64)  # Adjusted for python-pptx view
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "도구를 넘어선 시스템 구축으로 압도적 속도와 품질 산출하기"
    p2.font.size = Pt(28)
    p2.font.color.rgb = c_gold
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2: Objective
    slide = prs.slides.add_slide(blank_layout)
    # Title
    t_box = slide.shapes.add_textbox(safe_margin, Inches(0.8), Inches(11), Inches(1.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = "미팅 목적: AI 중심 프로세스 재설계"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = c_black

    b_box = slide.shapes.add_textbox(safe_margin, Inches(2.0), Inches(11), Inches(4.5))
    tf = b_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI를 '보조 도구'로 쓰는 수준을 넘어, 앱 제작의 전체 흐름을 AI 중심으로 재정의합니다."
    p.font.size = Pt(24)
    p.font.color.rgb = c_blue
    
    bullets = [
        "1. 극단적 속도: 사이클 단축 및 문서화 최소화",
        "2. 일관된 품질: 변동폭 없는 프로덕트급 결과물 유지",
        "3. 운영 최적화: 실제 배포 및 운영 가능한 수준의 산출물 연결"
    ]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = c_grey
        p.level = 1

    # Slide 3: Divider 1
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_blue
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(safe_margin, Inches(4), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "01. 기존 방식 vs AI 중심 방식"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.RIGHT

    # Slide 4: Traditional Flow
    slide = prs.slides.add_slide(blank_layout)
    t_box = slide.shapes.add_textbox(safe_margin, Inches(0.8), Inches(11), Inches(1.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = "기존 앱 제작 흐름 (전통적 방식)"
    p.font.size = Pt(44)
    p.font.bold = True

    b_box = slide.shapes.add_textbox(safe_margin, Inches(2.0), Inches(11), Inches(4))
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.text = "아이디어 도출 ➔ 기획 ➔ BM 검증 ➔ BI 작업 ➔ 와이어프레임 ➔ 프론트엔드 개발"
    p.font.size = Pt(24)
    p.font.color.rgb = c_grey
    
    p = tf.add_paragraph()
    p.text = "[Insight] 중간 단계의 '문서화 및 회의'로 인한 긴 리드 타임 발생."
    p.font.size = Pt(24)
    p.font.color.rgb = c_gold
    p.font.bold = True

    # Slide 5: AI-Centric Flow
    slide = prs.slides.add_slide(blank_layout)
    t_box = slide.shapes.add_textbox(safe_margin, Inches(0.8), Inches(11), Inches(1.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = "AI 중심 방식: 프롬프트 압축형 설계"
    p.font.size = Pt(44)
    p.font.bold = True

    b_box = slide.shapes.add_textbox(safe_margin, Inches(2.0), Inches(11), Inches(2))
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.text = "아이디어 도출 ➔ 기획 내용을 자연어로 명령하여 즉시 프론트엔드로 구현"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = c_blue

    # Badge Box
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, safe_margin, Inches(3.5), Inches(10), Inches(1.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(248, 248, 248)
    badge.line.fill.background()
    p = badge.text_frame.paragraphs[0]
    p.text = "중간 단계를 길게 끌지 않고, 프롬프트로 압축해 바로 산출물을 만든다."
    p.font.size = Pt(28)
    p.font.color.rgb = c_deep_slate
    p.alignment = PP_ALIGN.CENTER

    # Slide 6: Divider 2
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_purple
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(safe_margin, Inches(4), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "02. 품질의 핵심: 구체화와 시스템"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.RIGHT

    # Slide 8: Design System Def
    slide = prs.slides.add_slide(blank_layout)
    t_box = slide.shapes.add_textbox(safe_margin, Inches(0.8), Inches(11), Inches(1.5))
    p = t_box.text_frame.paragraphs[0]
    p.text = "프롬프트를 구체화시키는 핵심 도구 = 「디자인 시스템」"
    p.font.size = Pt(40)
    p.font.bold = True

    b_box = slide.shapes.add_textbox(safe_margin, Inches(2.0), Inches(11), Inches(4))
    tf = b_box.text_frame
    p = tf.paragraphs[0]
    p.text = "- 단순한 UI(색/폰트) 개념이 아닙니다."
    p.font.size = Pt(28)
    p = tf.add_paragraph()
    p.text = "- 앱이 실제로 동작하고 운영되는 수준까지 포함한 전체 시스템 설계."
    p.font.size = Pt(28)
    p.font.color.rgb = c_gold
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = "\n[Flow 적용]"
    p.font.size = Pt(24)
    p.font.color.rgb = c_blue
    p = tf.add_paragraph()
    p.text = "시스템 설계 ➔ 피그마메이크/안티그래비티/클로드코드 입력 ➔ 일관된 코드로 도출"
    p.font.size = Pt(24)

    # Slide 10: 8 Steps
    slide = prs.slides.add_slide(blank_layout)
    t_box = slide.shapes.add_textbox(safe_margin, Inches(0.5), Inches(11), Inches(1.0))
    p = t_box.text_frame.paragraphs[0]
    p.text = "프로덕트를 아우르는 확장성 (8단계 범위)"
    p.font.size = Pt(40)
    p.font.bold = True

    steps = [
        "[1] 문제 정의 & 사용자 구조", "[2] 인증·권한·상태 모델", 
        "[3] 기능 상태 & 데이터 표현 구조", "[4] 인터랙션 흐름 설계", 
        "[5] UI 토큰 및 컴포넌트 적용", "[6] 결제·백오피스 도메인 확장", 
        "[7] 기술 및 시스템 확장 설계", "[8] 운영·관측·성장 체계"
    ]
    for i, step_text in enumerate(steps):
        row = i // 4
        col = i % 4
        w = Inches(2.6)
        h = Inches(1.5)
        l = safe_margin + col * Inches(2.8)
        t = Inches(1.8) + row * Inches(1.8)
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = c_deep_slate
        p = box.text_frame.paragraphs[0]
        p.text = step_text
        p.font.size = Pt(20)
        p.font.color.rgb = c_white
        p.alignment = PP_ALIGN.CENTER

    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, safe_margin, Inches(5.8), Inches(11.1), Inches(1.0))
    badge.fill.solid()
    badge.fill.fore_color.rgb = c_gold
    p = badge.text_frame.paragraphs[0]
    p.text = "결론: 예쁜 화면이 아니라, 운영 체계를 전부 포함하는 개념."
    p.font.size = Pt(24)
    p.font.color.rgb = c_deep_slate
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Slide 11: Final Goal
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = c_deep_slate
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(safe_margin, Inches(2), Inches(11), Inches(4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "최종 목적지: 원본을 벗어난 '나만의 시스템'"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "\n컬러, 폰트, 토큰, 컴포넌트 룰, 모델 규칙을 지속적으로 변형.\n원본 데이터와 멀어지는 과정을 거쳐 가장 강력한 독자적 시스템을 완성한다."
    p.font.size = Pt(28)
    p.font.color.rgb = c_gold
    p.alignment = PP_ALIGN.CENTER

    prs.save("/Users/minwoo/Documents/New project/AI_Centric_Presentation.pptx")
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_pres()
